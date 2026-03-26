from __future__ import annotations

import re
import tempfile
from collections import Counter
import logging
from pathlib import Path
from typing import Iterable, List, Optional

from docx import Document as DocxDocument
from pypdf import PdfReader
from pptx import Presentation

from .models import CommitteeDocument, DocumentSection

logger = logging.getLogger(__name__)

_STAT_PATTERN = re.compile(
    r"(?:\$?\d[\d,]*(?:\.\d+)?\s?(?:[KMBkmb]|million|billion|trillion|hours|days|jobs|docs|states|%|x)?)"
)
_SENTENCE_SPLIT = re.compile(r"(?<=[.!?])\s+")


def parse_document_bytes(filename: str, payload: bytes) -> CommitteeDocument:
    warnings: List[str] = []
    parse_mode = "native"
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        sections, parse_mode, warnings = _parse_pdf(filename, payload)
        source_type = "pdf"
    elif ext == ".pptx":
        sections = _parse_pptx(payload)
        source_type = "pptx"
    elif ext == ".docx":
        sections = _parse_docx(payload)
        source_type = "docx"
    else:
        text = payload.decode("utf-8", errors="ignore")
        sections = _chunk_text_sections(text)
        source_type = "txt"

    raw_text = "\n\n".join(section.content for section in sections)
    return CommitteeDocument(
        documentId=_slugify(filename),
        documentName=filename,
        sourceType=source_type,
        sections=sections,
        rawText=raw_text,
        parseMode=parse_mode,  # type: ignore[arg-type]
        parseWarnings=warnings,
    )


def _parse_pdf(filename: str, payload: bytes) -> tuple[List[DocumentSection], str, List[str]]:
    sections: List[DocumentSection] = []
    warnings: List[str] = []
    try:
        with tempfile.NamedTemporaryFile(suffix=".pdf") as handle:
            handle.write(payload)
            handle.flush()
            reader = PdfReader(handle.name)
            raw_pages: List[List[str]] = []
            counter: Counter[str] = Counter()
            for page in reader.pages:
                lines = [_normalize_line(line) for line in (page.extract_text() or "").splitlines()]
                lines = [line for line in lines if line]
                raw_pages.append(lines)
                counter.update(set(lines))

            repeated_threshold = max(3, int(len(raw_pages) * 0.6))
            repeated_lines = {
                line for line, count in counter.items()
                if count >= repeated_threshold or _looks_like_pdf_chrome(line)
            }

            for index, lines in enumerate(raw_pages, start=1):
                cleaned_lines = [
                    line for line in lines
                    if line not in repeated_lines
                    and not _is_page_number(line)
                    and not _looks_like_slide_artifact(line)
                ]
                if cleaned_lines:
                    sections.append(_section_from_lines(index, cleaned_lines, [index]))
    except Exception as exc:
        logger.warning("Primary PDF parsing failed for %s: %s", filename, exc)
        warnings.append("Primary PDF parsing failed; retried in compatibility mode.")
        sections = []

    if sections:
        return sections, "native", warnings

    # Fallback: use pdftotext if extraction produced empty output.
    try:
        import subprocess

        with tempfile.NamedTemporaryFile(suffix=".pdf") as src, tempfile.NamedTemporaryFile(suffix=".txt") as out:
            src.write(payload)
            src.flush()
            subprocess.run(
                ["pdftotext", "-layout", src.name, out.name],
                check=True,
                capture_output=True,
                text=True,
            )
            text = Path(out.name).read_text(encoding="utf-8", errors="ignore")
        compat_sections = _chunk_text_sections(text)
        if compat_sections:
            warnings.append("PDF text was extracted in compatibility mode; section boundaries may be approximate.")
            return compat_sections, "compatibility", warnings
    except Exception as exc:
        logger.warning("Compatibility PDF parsing failed for %s: %s", filename, exc)
        warnings.append("Compatibility PDF parsing failed; used raw text fallback.")
        fallback = payload.decode("utf-8", errors="ignore")
        return _chunk_text_sections(fallback or filename), "fallback", warnings

    warnings.append("PDF produced no structured text; used raw fallback text.")
    fallback = payload.decode("utf-8", errors="ignore")
    return _chunk_text_sections(fallback or filename), "fallback", warnings


def _parse_docx(payload: bytes) -> List[DocumentSection]:
    with tempfile.NamedTemporaryFile(suffix=".docx") as handle:
        handle.write(payload)
        handle.flush()
        document = DocxDocument(handle.name)
        paragraphs = [paragraph for paragraph in document.paragraphs if paragraph.text.strip()]

    sections: List[DocumentSection] = []
    current_title: Optional[str] = None
    current_lines: List[str] = []

    for paragraph in paragraphs:
        text = paragraph.text.strip()
        style_name = getattr(paragraph.style, "name", "") or ""
        is_heading = style_name.lower().startswith("heading")
        if is_heading and current_lines:
            sections.append(
                _section_from_text(
                    len(sections) + 1,
                    "\n".join(current_lines),
                    [],
                    title=current_title,
                )
            )
            current_lines = []
        if is_heading:
            current_title = text
        else:
            if not current_title and len(text.split()) <= 12:
                current_title = text
            current_lines.append(text)

    if current_lines:
        sections.append(
            _section_from_text(
                len(sections) + 1,
                "\n".join(current_lines),
                [],
                title=current_title,
            )
        )

    return sections or _chunk_text_sections("")


def _parse_pptx(payload: bytes) -> List[DocumentSection]:
    with tempfile.NamedTemporaryFile(suffix=".pptx") as handle:
        handle.write(payload)
        handle.flush()
        presentation = Presentation(handle.name)
        sections: List[DocumentSection] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            texts: List[str] = []
            title: Optional[str] = None
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                value = shape.text.strip()
                if not value:
                    continue
                if title is None:
                    title = value.splitlines()[0].strip()
                texts.append(value)
            merged = "\n".join(texts).strip()
            if merged:
                sections.append(_section_from_text(slide_index, merged, [slide_index], title=title))
        return sections or _chunk_text_sections("")


def _chunk_text_sections(text: str) -> List[DocumentSection]:
    blocks = [chunk.strip() for chunk in re.split(r"\n\s*\n+", text) if chunk.strip()]
    if not blocks:
        blocks = [text.strip()] if text.strip() else ["Untitled document section"]

    sections: List[DocumentSection] = []
    for index, block in enumerate(blocks, start=1):
        sections.append(_section_from_text(index, block, []))
    return sections


def _section_from_text(
    index: int,
    text: str,
    slide_refs: List[int],
    title: Optional[str] = None,
) -> DocumentSection:
    normalized = re.sub(r"\s+\n", "\n", text).strip()
    lines = [line.strip(" -\u2022\t") for line in normalized.splitlines() if line.strip()]
    derived_title = title or (lines[0][:90] if lines else f"Section {index}")
    content_lines = lines[1:] if title is None and len(lines) > 1 else lines
    content = "\n".join(content_lines).strip() or normalized
    stats = _extract_stats(normalized)
    claims = _extract_claims(content)
    proof_points = _extract_proof_points(content)
    cta = _extract_cta(content)
    section_type = _infer_section_type(derived_title, content, stats, cta)

    return DocumentSection(
        id=index,
        title=derived_title,
        content=content,
        type=section_type,
        slideRefs=slide_refs,
        stats=stats,
        claims=claims,
        proofPoints=proof_points,
        cta=cta,
    )


def _section_from_lines(
    index: int,
    lines: List[str],
    slide_refs: List[int],
) -> DocumentSection:
    title = _select_title(lines, index)
    content_lines = [line for line in lines if line != title] if title else lines
    content = "\n".join(content_lines).strip() or "\n".join(lines).strip()
    return _section_from_text(index, content, slide_refs, title=title)


def _extract_stats(text: str) -> List[str]:
    values = []
    for match in _STAT_PATTERN.findall(text):
        value = match.strip()
        if value and any(ch.isdigit() for ch in value):
            values.append(value)
    return list(dict.fromkeys(values))[:8]


def _extract_claims(text: str) -> List[str]:
    sentences = [sentence.strip() for sentence in _SENTENCE_SPLIT.split(text) if sentence.strip()]
    candidates = [sentence for sentence in sentences if 6 <= len(sentence.split()) <= 24]
    return candidates[:4]


def _extract_proof_points(text: str) -> List[str]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    keywords = ("report", "hours", "saved", "faster", "federal", "fedramp", "oig", "cisa", "search.gov")
    matches = [line for line in lines if any(keyword in line.lower() for keyword in keywords)]
    return matches[:4]


def _extract_cta(text: str) -> Optional[str]:
    for sentence in _SENTENCE_SPLIT.split(text):
        lowered = sentence.lower()
        if any(token in lowered for token in ("next step", "schedule", "discovery", "demo", "call")):
            return sentence.strip()
    return None


def _normalize_line(line: str) -> str:
    normalized = line.replace("\u00a0", " ").replace("ʼ", "'").replace("→", " to ").replace("➔", "- ")
    normalized = normalized.replace("\ue0a3", "").replace("\ue088", "")
    normalized = re.sub(r"\s+", " ", normalized)
    return normalized.strip(" |")


def _looks_like_pdf_chrome(line: str) -> bool:
    lowered = line.lower()
    return any(token in lowered for token in ("elastic perspective", "elastic.co", "all rights reserved"))


def _is_page_number(line: str) -> bool:
    return bool(re.fullmatch(r"\d{1,2}", line.strip()))


def _looks_like_slide_artifact(line: str) -> bool:
    if len(line) <= 2 and line.isupper():
        return True
    return False


def _select_title(lines: List[str], index: int) -> str:
    if not lines:
        return f"Section {index}"

    if index == 1:
        title_parts: List[str] = []
        for line in lines[:5]:
            if _is_metadata_line(line):
                break
            if len(line.split()) <= 8:
                title_parts.append(line)
            else:
                break
        if len(title_parts) >= 2:
            return " ".join(title_parts)

    candidates = lines[:12]
    best = candidates[0]
    best_score = -10.0
    for line in candidates:
        score = _heading_score(line)
        if score > best_score:
            best = line
            best_score = score
    return best


def _is_metadata_line(line: str) -> bool:
    lowered = line.lower()
    return any(token in lowered for token in ("march", "built for", "an elastic perspective", "sources:", "copyright", "©", "account team"))


def _heading_score(line: str) -> float:
    words = line.split()
    if not words:
        return -10.0
    score = 0.0
    if 2 <= len(words) <= 12:
        score += 2.0
    if len(line) <= 80:
        score += 1.0
    if any(char.isdigit() for char in line):
        score -= 0.8
    if line.endswith("."):
        score -= 0.4
    title_case_ratio = sum(1 for word in words if word[:1].isupper()) / max(len(words), 1)
    score += title_case_ratio
    lowered = line.lower()
    if lowered == "agenda":
        return 10.0
    if any(token in lowered for token in ("agenda", "summary", "priorities", "search", "proof", "risk", "action", "opportunity", "overview", "analysis", "landscape", "deployment", "impact", "inconsistency")):
        score += 1.5
    if ":" in line:
        score += 0.3
    return score


def _infer_section_type(title: str, content: str, stats: Iterable[str], cta: Optional[str]) -> str:
    lowered_title = title.lower()
    lowered_content = content.lower()
    if "appendix" in lowered_title:
        return "appendix"
    if "thank you" in lowered_title:
        return "closing"
    if any(token in lowered_title for token in ("agenda", "outline")):
        return "agenda"
    if cta:
        return "call_to_action"
    if stats:
        return "evidence"
    if any(token in lowered_content for token in ("case study", "search.gov", "cisa", "docusign")):
        return "proof"
    if len(content.split()) < 20:
        return "title"
    return "content"


def _slugify(value: str) -> str:
    base = Path(value).stem.lower()
    base = re.sub(r"[^a-z0-9]+", "-", base).strip("-")
    return base or "committee-document"
